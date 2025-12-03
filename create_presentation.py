from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def set_text_frame_properties(tf, font_size=18):
    """Configure text frame to auto-fit and set font size"""
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.3)
    
    # Set font size for all paragraphs
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
        if paragraph.level == 1:
            for run in paragraph.runs:
                run.font.size = Pt(font_size - 2)
        elif paragraph.level > 1:
            for run in paragraph.runs:
                run.font.size = Pt(font_size - 4)

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Healthcare Provider Fraud Detection System"
subtitle.text = "(Medicare Fraud Detection Capstone)\n\nTeam: Mohamed Osama • Malak Khaled • Bahy Hany • Mohamed Wael"

# Set font sizes for title slide
title.text_frame.paragraphs[0].font.size = Pt(44)
if len(subtitle.text_frame.paragraphs) > 0:
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)

# Slide 2: The Problem & Objective
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
content2 = slide2.placeholders[1]

title2.text = "Slide 1 — The Problem & Our Objective\nBy Mohamed"
title2.text_frame.paragraphs[0].font.size = Pt(32)

tf2 = content2.text_frame
tf2.text = "US healthcare fraud exceeds $68B/year"
p = tf2.add_paragraph()
p.text = "Legacy auditing misses hidden fraud patterns"
p = tf2.add_paragraph()
p.text = "Fraudulent providers ≈ 9% → ML bias"
p = tf2.add_paragraph()
p.text = ""
p = tf2.add_paragraph()
p.text = "Goal"
p.level = 0
p = tf2.add_paragraph()
p.text = "Detect high-risk providers (not individual claims)"
p.level = 1
p = tf2.add_paragraph()
p.text = "Handle severe class imbalance"
p.level = 1
p = tf2.add_paragraph()
p.text = "Reduce false investigations & financial loss"
p.level = 1
set_text_frame_properties(tf2, font_size=18)

# Slide 3: Data Architecture
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
content3 = slide3.placeholders[1]

title3.text = "Slide 2 — Data Architecture Overview (Phase 1)\nBy Malak"
title3.text_frame.paragraphs[0].font.size = Pt(28)

tf3 = content3.text_frame
tf3.text = "Inputs"
p = tf3.add_paragraph()
p.text = "Beneficiary"
p.level = 1
p = tf3.add_paragraph()
p.text = "Inpatient"
p.level = 1
p = tf3.add_paragraph()
p.text = "Outpatient"
p.level = 1
p = tf3.add_paragraph()
p.text = "Provider labels"
p.level = 1
p = tf3.add_paragraph()
p.text = ""
p = tf3.add_paragraph()
p.text = "Join Level"
p.level = 0
p = tf3.add_paragraph()
p.text = "BeneID → Patient Events"
p.level = 1
p = tf3.add_paragraph()
p.text = "Provider → Target Entity (Fraud label)"
p.level = 1
p = tf3.add_paragraph()
p.text = ""
p = tf3.add_paragraph()
p.text = "Challenge"
p.level = 0
p = tf3.add_paragraph()
p.text = "Claims are transactional"
p.level = 1
p = tf3.add_paragraph()
p.text = "Fraud is behavioral at provider level"
p.level = 1
p = tf3.add_paragraph()
p.text = ""
p = tf3.add_paragraph()
p.text = "➡️ Build provider-level profiles"
set_text_frame_properties(tf3, font_size=18)

# Slide 4: Data Cleaning
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
content4 = slide4.placeholders[1]

title4.text = "Slide 3 — Data Cleaning Strategy\nBy Malak"
title4.text_frame.paragraphs[0].font.size = Pt(32)

tf4 = content4.text_frame
tf4.text = "Problems Detected"
p = tf4.add_paragraph()
p.text = "Missing physician codes"
p.level = 1
p = tf4.add_paragraph()
p.text = "Nonsensical claim durations"
p.level = 1
p = tf4.add_paragraph()
p.text = "Zero-information columns"
p.level = 1
p = tf4.add_paragraph()
p.text = ""
p = tf4.add_paragraph()
p.text = "Actions"
p.level = 0
p = tf4.add_paragraph()
p.text = "Encode missing physicians as \"Unknown\""
p.level = 1
p = tf4.add_paragraph()
p.text = "Drop invalid-duration claims"
p.level = 1
p = tf4.add_paragraph()
p.text = "Remove redundant fields"
p.level = 1
p = tf4.add_paragraph()
p.text = ""
p = tf4.add_paragraph()
p.text = "Outcome"
p.level = 0
p = tf4.add_paragraph()
p.text = "Fraud distribution confirmed at ~9–10%"
p.level = 1
set_text_frame_properties(tf4, font_size=18)

# Slide 5: Feature Engineering
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
content5 = slide5.placeholders[1]

title5.text = "Slide 4 — Feature Engineering\nBy Mohamed"
title5.text_frame.paragraphs[0].font.size = Pt(32)

tf5 = content5.text_frame
tf5.text = "We convert millions of claims → provider behavioral vectors"
p = tf5.add_paragraph()
p.text = ""
p = tf5.add_paragraph()
p.text = "Feature Groups"
p.level = 0
p = tf5.add_paragraph()
p.text = "Financial: reimbursements, avg deductible, mean claim cost"
p.level = 1
p = tf5.add_paragraph()
p.text = "Operational: volume, inpatient/outpatient balance"
p.level = 1
p = tf5.add_paragraph()
p.text = "Service Diversity: unique diagnoses & procedures"
p.level = 1
p = tf5.add_paragraph()
p.text = "Patient mix: chronic ratios, unique patient count"
p.level = 1
p = tf5.add_paragraph()
p.text = ""
p = tf5.add_paragraph()
p.text = "Final Dataset: 5410 providers × ~134 features"
set_text_frame_properties(tf5, font_size=18)

# Slide 6: Exploratory Insights
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title6 = slide6.shapes.title
content6 = slide6.placeholders[1]

title6.text = "Slide 5 — Exploratory Insights\nBy Mohamed"
title6.text_frame.paragraphs[0].font.size = Pt(32)

tf6 = content6.text_frame
tf6.text = "Detected patterns:"
p = tf6.add_paragraph()
p.text = "Fraud providers show higher revenue throughput"
p.level = 1
p = tf6.add_paragraph()
p.text = "Greater diagnosis/procedure diversity"
p.level = 1
p = tf6.add_paragraph()
p.text = "Geographic hotspots (certain states)"
p.level = 1
p = tf6.add_paragraph()
p.text = "Sustained high billing across time → not random"
p.level = 1
set_text_frame_properties(tf6, font_size=18)

# Slide 7: Handling Imbalance
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
title7 = slide7.shapes.title
content7 = slide7.placeholders[1]

title7.text = "Slide 6 — Handling Imbalance (Phase 2)\nBy Bahy"
title7.text_frame.paragraphs[0].font.size = Pt(28)

tf7 = content7.text_frame
tf7.text = "Issue: Fraud = minority → models default to \"No Fraud\""
p = tf7.add_paragraph()
p.text = ""
p = tf7.add_paragraph()
p.text = "Methods tested:"
p.level = 0
p = tf7.add_paragraph()
p.text = "Class Weighting"
p.level = 1
p = tf7.add_paragraph()
p.text = "SMOTE"
p.level = 1
p = tf7.add_paragraph()
p.text = "SMOTE + Tomek Links"
p.level = 1
p = tf7.add_paragraph()
p.text = ""
p = tf7.add_paragraph()
p.text = "Purpose:"
p.level = 0
p = tf7.add_paragraph()
p.text = "Improve detection power"
p.level = 1
p = tf7.add_paragraph()
p.text = "Preserve real distributions after oversampling"
p.level = 1
set_text_frame_properties(tf7, font_size=18)

# Slide 8: Modeling Strategy
slide8 = prs.slides.add_slide(prs.slide_layouts[1])
title8 = slide8.shapes.title
content8 = slide8.placeholders[1]

title8.text = "Slide 7 — Modeling Strategy (Phase 3)\nBy Bahy"
title8.text_frame.paragraphs[0].font.size = Pt(28)

tf8 = content8.text_frame
tf8.text = "Models:"
p = tf8.add_paragraph()
p.text = "Logistic Regression → interpretable baseline"
p.level = 1
p = tf8.add_paragraph()
p.text = "Decision Tree → simple benchmark"
p.level = 1
p = tf8.add_paragraph()
p.text = "Random Forest → robust ensemble"
p.level = 1
p = tf8.add_paragraph()
p.text = "XGBoost → non-linear excellence + imbalance tolerance"
p.level = 1
p = tf8.add_paragraph()
p.text = ""
p = tf8.add_paragraph()
p.text = "Tuning"
p.level = 0
p = tf8.add_paragraph()
p.text = "GridSearchCV (5-fold)"
p.level = 1
p = tf8.add_paragraph()
p.text = "Optimized for F1 score"
p.level = 1
set_text_frame_properties(tf8, font_size=18)

# Slide 9: Model Selection Results
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
title9 = slide9.shapes.title
content9 = slide9.placeholders[1]

title9.text = "Slide 8 — Model Selection Results\nBy Mohamed"
title9.text_frame.paragraphs[0].font.size = Pt(28)

tf9 = content9.text_frame
tf9.text = "Model Comparison (Validation Set):"
p = tf9.add_paragraph()
p.text = ""
p = tf9.add_paragraph()
p.text = "Logistic Regression: Precision 0.43, Recall 0.80, F1 0.56, PR-AUC 0.63"
p = tf9.add_paragraph()
p.text = "Random Forest: Precision 0.48, Recall 0.77, F1 0.59, PR-AUC 0.59"
p = tf9.add_paragraph()
p.text = "Decision Tree: Precision 0.43, Recall 0.68, F1 0.53, PR-AUC 0.37"
p = tf9.add_paragraph()
p.text = "XGBoost: Precision 0.59, Recall 0.68, F1 0.63, PR-AUC 0.65"
p = tf9.add_paragraph()
p.text = ""
p = tf9.add_paragraph()
p.text = "XGBoost = best balance: recall + precision + ROC-AUC"
set_text_frame_properties(tf9, font_size=16)

# Slide 10: Final Test Evaluation
slide10 = prs.slides.add_slide(prs.slide_layouts[1])
title10 = slide10.shapes.title
content10 = slide10.placeholders[1]

title10.text = "Slide 9 — Final Test Evaluation (Phase 4)\nBy Aly"
title10.text_frame.paragraphs[0].font.size = Pt(28)

tf10 = content10.text_frame
tf10.text = "Test Results (20% unseen)"
p = tf10.add_paragraph()
p.text = "Accuracy: 93.63%"
p.level = 1
p = tf10.add_paragraph()
p.text = "Precision: 63.79%"
p.level = 1
p = tf10.add_paragraph()
p.text = "Recall: 73.27%"
p.level = 1
p = tf10.add_paragraph()
p.text = "F1: 68.20%"
p.level = 1
p = tf10.add_paragraph()
p.text = "ROC-AUC: 96.79%"
p.level = 1
p = tf10.add_paragraph()
p.text = "PR-AUC: 79.89%"
p.level = 1
p = tf10.add_paragraph()
p.text = ""
p = tf10.add_paragraph()
p.text = "Confusion Matrix"
p.level = 0
p = tf10.add_paragraph()
p.text = "TP: 74 — caught fraud"
p.level = 1
p = tf10.add_paragraph()
p.text = "FN: 27 — missed fraud"
p.level = 1
p = tf10.add_paragraph()
p.text = "FP: 42 — unnecessary audits"
p.level = 1
p = tf10.add_paragraph()
p.text = "TN: 940 — legitimate providers"
p.level = 1
p = tf10.add_paragraph()
p.text = ""
p = tf10.add_paragraph()
p.text = "➡️ High recall + very high ROC-AUC = excellent fraud ranking"
set_text_frame_properties(tf10, font_size=16)

# Slide 11: Business Impact
slide11 = prs.slides.add_slide(prs.slide_layouts[1])
title11 = slide11.shapes.title
content11 = slide11.placeholders[1]

title11.text = "Slide 10 — Business Impact\nBy Mohamed"
title11.text_frame.paragraphs[0].font.size = Pt(32)

tf11 = content11.text_frame
tf11.text = "Assumptions:"
p = tf11.add_paragraph()
p.text = "False Positive → ~$1,000 audit cost"
p.level = 1
p = tf11.add_paragraph()
p.text = "False Negative → ~$50,000 undetected fraud"
p.level = 1
p = tf11.add_paragraph()
p.text = ""
p = tf11.add_paragraph()
p.text = "Model advantages:"
p.level = 0
p = tf11.add_paragraph()
p.text = "Prevents silent financial leakage"
p.level = 1
p = tf11.add_paragraph()
p.text = "Strong ROC-AUC → optimal investigation prioritization"
p.level = 1
p = tf11.add_paragraph()
p.text = "Saves operational effort vs random review"
p.level = 1
set_text_frame_properties(tf11, font_size=18)

# Slide 12: Error Analysis
slide12 = prs.slides.add_slide(prs.slide_layouts[1])
title12 = slide12.shapes.title
content12 = slide12.placeholders[1]

title12.text = "Slide 11 — Error Analysis\nBy Bahy"
title12.text_frame.paragraphs[0].font.size = Pt(32)

tf12 = content12.text_frame
tf12.text = "False Positives"
p = tf12.add_paragraph()
p.text = "Legit high-throughput providers"
p.level = 1
p = tf12.add_paragraph()
p.text = "Resemble \"high-revenue fraud\""
p.level = 1
p = tf12.add_paragraph()
p.text = ""
p = tf12.add_paragraph()
p.text = "False Negatives"
p.level = 0
p = tf12.add_paragraph()
p.text = "Very low-volume clinics"
p.level = 1
p = tf12.add_paragraph()
p.text = "Mimic normal small providers"
p.level = 1
p = tf12.add_paragraph()
p.text = ""
p = tf12.add_paragraph()
p.text = "Future Enhancements"
p.level = 0
p = tf12.add_paragraph()
p.text = "Provider–patient graph relationships"
p.level = 1
p = tf12.add_paragraph()
p.text = "Temporal fraud drift tracking"
p.level = 1
set_text_frame_properties(tf12, font_size=18)

# Slide 13: Key Drivers
slide13 = prs.slides.add_slide(prs.slide_layouts[1])
title13 = slide13.shapes.title
content13 = slide13.placeholders[1]

title13.text = "Slide 12 — Key Drivers of Risk\nBy Malak"
title13.text_frame.paragraphs[0].font.size = Pt(32)

tf13 = content13.text_frame
tf13.text = "Top influential signals:"
p = tf13.add_paragraph()
p.text = "Total claims submitted"
p.level = 1
p = tf13.add_paragraph()
p.text = "Total & avg reimbursement"
p.level = 1
p = tf13.add_paragraph()
p.text = "Diagnosis & procedure diversity"
p.level = 1
p = tf13.add_paragraph()
p.text = "Chronic illness ratios"
p.level = 1
p = tf13.add_paragraph()
p.text = ""
p = tf13.add_paragraph()
p.text = "➡️ Features align with known real-world fraud behaviors"
set_text_frame_properties(tf13, font_size=18)

# Slide 14: Recommendations
slide14 = prs.slides.add_slide(prs.slide_layouts[1])
title14 = slide14.shapes.title
content14 = slide14.placeholders[1]

title14.text = "Slide 13 — Final Recommendations\nBy Mohamed"
title14.text_frame.paragraphs[0].font.size = Pt(28)

tf14 = content14.text_frame
tf14.text = "Deploy XGBoost for Medicare fraud triage"
p = tf14.add_paragraph()
p.text = "Investigate top risk percentile providers"
p = tf14.add_paragraph()
p.text = "Keep human analysts in loop"
p = tf14.add_paragraph()
p.text = "Retrain every quarter → prevent drift"
set_text_frame_properties(tf14, font_size=20)

# Slide 15: Q&A
slide15 = prs.slides.add_slide(prs.slide_layouts[0])
title15 = slide15.shapes.title
title15.text = "Slide 14 — Q&A\n\nQuestions?"
title15.text_frame.paragraphs[0].font.size = Pt(44)

# Save presentation
prs.save('reports/presentation.pptx')
print("PowerPoint presentation created successfully at reports/presentation.pptx")
print("All text has been configured to fit within slide boundaries with appropriate font sizes.")
